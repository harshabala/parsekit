#!/usr/bin/env python3
"""
ParseKit token benchmark: naive extraction vs sidecar Markdown/JSON output.

Install deps (once):
  python3 -m venv .venv-benchmark
  .venv-benchmark/bin/pip install -r scripts/requirements-benchmark.txt

Run:
  python3 scripts/benchmark_tokens.py
"""

from __future__ import annotations

import json
import os
import shutil
import subprocess
import sys
import tempfile
import zipfile
from dataclasses import dataclass
from datetime import datetime, timezone
from pathlib import Path
from typing import Callable
from xml.etree import ElementTree as ET

ROOT = Path(__file__).resolve().parent.parent
FIXTURES_DIR = ROOT / "scripts" / "benchmark-fixtures"
RESULTS_MD = ROOT / "docs" / "benchmark-results.md"
RATIOS_JSON = ROOT / "scripts" / "benchmark-ratios.json"
REQUIREMENTS = ROOT / "scripts" / "requirements-benchmark.txt"

# Office Open XML namespaces used for naive baseline extraction.
OOXML_NS = {
    "w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main",
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "main": "http://schemas.openxmlformats.org/spreadsheetml/2006/main",
}


@dataclass
class FixtureSpec:
    filename: str
    file_type: str
    kind: str
    why: str


FIXTURES: list[FixtureSpec] = [
    FixtureSpec(
        "born-digital.pdf",
        "pdf",
        "born-digital",
        "Born-digital PDF text layer vs structured Markdown headings.",
    ),
    FixtureSpec(
        "scanned.pdf",
        "pdf",
        "scanned",
        "Image-only PDF: naive extract is empty; ParseKit OCR unlocks text (not a token reduction).",
    ),
    FixtureSpec(
        "sample.docx",
        "docx",
        "office",
        "Office XML overhead vs clean Markdown paragraphs.",
    ),
    FixtureSpec(
        "slides.pptx",
        "pptx",
        "office",
        "Deck slide/master XML vs Markdown slide sections.",
    ),
    FixtureSpec(
        "spreadsheet.xlsx",
        "xlsx",
        "office",
        "Spreadsheet cell dump vs LiteParse JSON export (tabular formats always emit JSON).",
    ),
]


def _require_deps() -> None:
    try:
        import tiktoken  # noqa: F401
        return
    except ImportError:
        pass

    venv_python = ROOT / ".venv-benchmark" / "bin" / "python3"
    if venv_python.is_file() and os.environ.get("PARSEKIT_BENCHMARK_VENV") != "1":
        env = os.environ.copy()
        env["PARSEKIT_BENCHMARK_VENV"] = "1"
        proc = subprocess.run([str(venv_python), *sys.argv], env=env)
        raise SystemExit(proc.returncode)

    print(
        "Missing Python dependencies for benchmark.\n"
        "Setup:\n"
        "  python3 -m venv .venv-benchmark\n"
        f"  .venv-benchmark/bin/pip install -r {REQUIREMENTS.relative_to(ROOT)}",
        file=sys.stderr,
    )
    sys.exit(1)


def _host_triple() -> str:
    try:
        out = subprocess.run(
            ["rustc", "-Vv"],
            check=True,
            capture_output=True,
            text=True,
        )
        for line in out.stdout.splitlines():
            if line.startswith("host:"):
                return line.split(":", 1)[1].strip()
    except (FileNotFoundError, subprocess.CalledProcessError):
        pass
    import platform

    system = platform.system().lower()
    machine = platform.machine().lower()
    if system == "darwin" and machine in ("arm64", "aarch64"):
        return "aarch64-apple-darwin"
    if system == "darwin":
        return "x86_64-apple-darwin"
    if system == "linux":
        return f"{machine}-unknown-linux-gnu"
    if system == "windows":
        return f"{machine}-pc-windows-msvc"
    return f"{machine}-unknown"


def find_sidecar() -> Path:
    binaries = ROOT / "src-tauri" / "binaries"
    candidates = [
        binaries / f"parsekit-sidecar-{_host_triple()}",
        binaries / "parsekit-sidecar",
    ]
    for path in candidates:
        if path.is_file():
            return path
    raise FileNotFoundError(
        "parsekit-sidecar binary not found. Build with: npm run build:sidecar"
    )


def _write_born_digital_pdf(path: Path) -> None:
    from reportlab.lib.pagesizes import letter
    from reportlab.pdfgen import canvas

    c = canvas.Canvas(str(path), pagesize=letter)
    text = c.beginText(72, 720)
    for line in [
        "ParseKit Benchmark — Born-Digital PDF",
        "",
        "This synthetic document tests naive PDF text extraction versus",
        "structured Markdown output. Revenue Q1: $1,240,000. Headcount: 42.",
        "Key initiatives: platform migration, OCR pipeline, token savings tracker.",
    ]:
        text.textLine(line)
    c.drawText(text)
    c.save()


def _write_scanned_pdf(path: Path) -> None:
    import io

    from PIL import Image, ImageDraw
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.utils import ImageReader
    from reportlab.pdfgen import canvas

    img = Image.new("RGB", (620, 220), "white")
    draw = ImageDraw.Draw(img)
    draw.text(
        (24, 90),
        "ParseKit scanned PDF benchmark. Invoice #BENCH-2048. Total: $512.00.",
        fill="black",
    )
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)

    c = canvas.Canvas(str(path), pagesize=letter)
    c.drawImage(ImageReader(buf), 72, 500, width=468, height=166)
    c.save()


def _write_docx(path: Path) -> None:
    from docx import Document

    doc = Document()
    doc.add_heading("ParseKit Benchmark DOCX", 0)
    doc.add_paragraph(
        "Synthetic Word document for token counting baseline versus Markdown export."
    )
    doc.add_paragraph("Table row: Product Alpha — units 1200 — margin 18%.")
    doc.save(path)


def _write_pptx(path: Path) -> None:
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "ParseKit Benchmark PPTX"
    slide.placeholders[1].text = (
        "Slide notes: compare deck XML bloat to clean Markdown bullets."
    )
    prs.save(path)


def _write_xlsx(path: Path) -> None:
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "Benchmark"
    ws.append(["Region", "Revenue", "Units"])
    ws.append(["North", 420000, 880])
    ws.append(["South", 310000, 640])
    wb.save(path)


GENERATORS: dict[str, Callable[[Path], None]] = {
    "born-digital.pdf": _write_born_digital_pdf,
    "scanned.pdf": _write_scanned_pdf,
    "sample.docx": _write_docx,
    "slides.pptx": _write_pptx,
    "spreadsheet.xlsx": _write_xlsx,
}


def ensure_fixtures() -> list[Path]:
    FIXTURES_DIR.mkdir(parents=True, exist_ok=True)
    paths: list[Path] = []
    for spec in FIXTURES:
        out = FIXTURES_DIR / spec.filename
        if not out.is_file():
            GENERATORS[spec.filename](out)
        paths.append(out)
    return paths


def _pdftotext(path: Path) -> str | None:
    if not shutil.which("pdftotext"):
        return None
    proc = subprocess.run(
        ["pdftotext", "-layout", str(path), "-"],
        capture_output=True,
        text=True,
    )
    if proc.returncode != 0:
        return None
    return proc.stdout


def _pypdf_text(path: Path) -> str:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    chunks = []
    for page in reader.pages:
        chunks.append(page.extract_text() or "")
    return "\n".join(chunks)


def baseline_pdf(path: Path) -> tuple[str, str]:
    text = _pdftotext(path)
    if text is not None:
        return text, "pdftotext"
    return _pypdf_text(path), "pypdf text layer"


def _ooxml_text_nodes(xml_bytes: bytes, local_tags: set[str]) -> str:
    root = ET.fromstring(xml_bytes)
    parts: list[str] = []

    def walk(elem: ET.Element) -> None:
        tag = elem.tag.rsplit("}", 1)[-1]
        if tag in local_tags and elem.text:
            parts.append(elem.text)
        if elem.tail:
            parts.append(elem.tail)
        for child in elem:
            walk(child)

    walk(root)
    return " ".join(parts)


def baseline_docx(path: Path) -> tuple[str, str]:
    with zipfile.ZipFile(path) as zf:
        xml = zf.read("word/document.xml")
    text = _ooxml_text_nodes(xml, {"t"})
    return text, "raw DOCX document.xml text nodes"


def baseline_pptx(path: Path) -> tuple[str, str]:
    chunks: list[str] = []
    with zipfile.ZipFile(path) as zf:
        slide_names = sorted(
            n for n in zf.namelist() if n.startswith("ppt/slides/slide") and n.endswith(".xml")
        )
        for name in slide_names:
            chunks.append(_ooxml_text_nodes(zf.read(name), {"t"}))
    return "\n".join(chunks), "raw PPTX slide XML text nodes"


def _col_letters(cell_ref: str) -> str:
    return "".join(ch for ch in cell_ref if ch.isalpha())


def baseline_xlsx(path: Path) -> tuple[str, str]:
    shared: list[str] = []
    rows: dict[int, dict[str, str]] = {}

    with zipfile.ZipFile(path) as zf:
        if "xl/sharedStrings.xml" in zf.namelist():
            sst = ET.fromstring(zf.read("xl/sharedStrings.xml"))
            for si in sst.findall("main:si", OOXML_NS):
                texts = [t.text or "" for t in si.findall(".//main:t", OOXML_NS)]
                shared.append("".join(texts))

        sheet = ET.fromstring(zf.read("xl/worksheets/sheet1.xml"))
        for cell in sheet.findall(".//main:c", OOXML_NS):
            ref = cell.attrib.get("r", "")
            if not ref:
                continue
            col = _col_letters(ref)
            row_num = int("".join(ch for ch in ref if ch.isdigit()) or "0")
            value = ""
            cell_type = cell.attrib.get("t")
            v = cell.find("main:v", OOXML_NS)
            if v is not None and v.text is not None:
                if cell_type == "s":
                    value = shared[int(v.text)]
                else:
                    value = v.text
            is_elem = cell.find("main:is", OOXML_NS)
            if is_elem is not None:
                texts = [t.text or "" for t in is_elem.findall(".//main:t", OOXML_NS)]
                value = "".join(texts)
            rows.setdefault(row_num, {})[col] = value

    lines: list[str] = []
    for row_num in sorted(rows):
        cols = rows[row_num]
        ordered = [cols[k] for k in sorted(cols)]
        lines.append("\t".join(ordered))
    return "\n".join(lines), "raw XLSX cell values (TSV)"


BASELINE_EXTRACTORS: dict[str, Callable[[Path], tuple[str, str]]] = {
    "pdf": baseline_pdf,
    "docx": baseline_docx,
    "pptx": baseline_pptx,
    "xlsx": baseline_xlsx,
}


def count_tokens(text: str, encoding_name: str) -> int:
    import tiktoken

    enc = tiktoken.get_encoding(encoding_name)
    return len(enc.encode(text or ""))


def reduction_ratio(baseline: int, parsekit: int) -> float:
    if baseline <= 0:
        return 0.0
    return max(0.0, (baseline - parsekit) / baseline)


def run_sidecar(sidecar: Path, files: list[Path], output_dir: Path) -> dict[str, Path]:
    if output_dir.exists():
        shutil.rmtree(output_dir)
    output_dir.mkdir(parents=True)

    config = {
        "files": [str(p.resolve()) for p in files],
        "outputDir": str(output_dir.resolve()),
        "format": "md",
        "ocrEnabled": True,
        "ocrLanguage": "eng",
        "workers": 2,
    }
    proc = subprocess.run(
        [str(sidecar)],
        input=json.dumps(config) + "\n",
        text=True,
        capture_output=True,
        timeout=600,
    )
    if proc.returncode != 0:
        print(proc.stdout, file=sys.stderr)
        print(proc.stderr, file=sys.stderr)
        raise RuntimeError(f"sidecar exited {proc.returncode}")

    outputs: dict[str, Path] = {}
    errors: list[str] = []
    for line in proc.stdout.splitlines():
        line = line.strip()
        if not line:
            continue
        event = json.loads(line)
        if event.get("type") == "progress" and event.get("status") == "completed":
            source = event.get("sourcePath", "")
            out_path = event.get("path")
            if source and out_path:
                outputs[source] = Path(out_path)
        if event.get("type") == "progress" and event.get("status") == "error":
            errors.append(
                f"{event.get('sourcePath')}: {event.get('error', 'unknown error')}"
            )
        if event.get("type") == "done" and event.get("errors", 0) > 0 and errors:
            raise RuntimeError("sidecar reported errors:\n" + "\n".join(errors))

    if not outputs:
        raise RuntimeError(
            "sidecar produced no completed outputs; stdout:\n" + proc.stdout[:4000]
        )
    return outputs


def pct(ratio: float) -> str:
    return f"{ratio * 100:.1f}%"


def write_results(
    rows: list[dict],
    sidecar_path: Path,
    generated_at: str,
) -> None:
    RESULTS_MD.parent.mkdir(parents=True, exist_ok=True)

    lines = [
        "# ParseKit Token Benchmark Results",
        "",
        f"> Generated by `scripts/benchmark_tokens.py` on {generated_at}.",
        "> Token counts use OpenAI reference tokenizers `cl100k_base` (GPT-4 family) "
        "and `o200k_base` (GPT-4o family). These are reference encodings only — "
        "actual model tokenization may differ.",
        "",
        "## Results",
        "",
        "| Fixture | Type | Baseline method | Baseline cl100k | ParseKit cl100k | "
        "Δ cl100k | Baseline o200k | ParseKit o200k | Δ o200k | Why |",
        "| --- | --- | --- | ---: | ---: | ---: | ---: | ---: | ---: | --- |",
    ]

    for row in rows:
        lines.append(
            f"| {row['fixture']} | {row['file_type']} | {row['baseline_method']} | "
            f"{row['baseline_cl100k']} | {row['parsekit_cl100k']} | {row['delta_cl100k']} | "
            f"{row['baseline_o200k']} | {row['parsekit_o200k']} | {row['delta_o200k']} | "
            f"{row['why']} |"
        )

    lines.extend(
        [
            "",
            "## Methodology",
            "",
            "- **Fixtures:** Synthetic, license-safe samples under `scripts/benchmark-fixtures/`.",
            "- **Baseline:** Naive extraction — `pdftotext` when available, otherwise PDF text "
            "layer via `pypdf`; Office formats use raw OOXML text/cell dumps.",
            f"- **ParseKit:** `{sidecar_path.relative_to(ROOT)}` sidecar (`format: md`, "
            "`ocrEnabled: true`). Spreadsheets emit JSON per ParseKit rules.",
            "- **Tokenizers:** `tiktoken` encodings `cl100k_base` and `o200k_base`.",
            "- **Δ (delta):** Percent reduction from baseline to ParseKit output; floored at 0% "
            "when baseline is empty (scanned PDFs) or ParseKit output is larger.",
            "",
            "## Reproduce",
            "",
            "```bash",
            "python3 -m venv .venv-benchmark",
            ".venv-benchmark/bin/pip install -r scripts/requirements-benchmark.txt",
            "python3 scripts/benchmark_tokens.py",
            "```",
            "",
        ]
    )

    RESULTS_MD.write_text("\n".join(lines), encoding="utf-8")


def write_ratios(rows: list[dict], generated_at: str) -> None:
    by_type: dict[str, list[float]] = {}
    for row in rows:
        by_type.setdefault(row["file_type"], []).append(row["reduction_ratio_cl100k"])

    payload = {
        "generated_at": generated_at,
        "source": "scripts/benchmark_tokens.py",
        "tokenizer_primary": "cl100k_base",
        "note": (
            "avg_reduction_ratio uses cl100k_base counts; scanned PDFs with zero baseline "
            "tokens contribute 0.0. Used by Task 5 token tracker for estimation fallback."
        ),
        "by_file_type": {
            file_type: {
                "files": len(ratios),
                "avg_reduction_ratio": round(sum(ratios) / len(ratios), 6),
            }
            for file_type, ratios in sorted(by_type.items())
        },
    }
    RATIOS_JSON.write_text(json.dumps(payload, indent=2) + "\n", encoding="utf-8")


def main() -> int:
    _require_deps()
    generated_at = datetime.now(timezone.utc).strftime("%Y-%m-%d %H:%M UTC")

    fixture_paths = ensure_fixtures()
    sidecar = find_sidecar()

    rows: list[dict] = []
    with tempfile.TemporaryDirectory(prefix="parsekit-bench-") as tmp:
        output_dir = Path(tmp)
        outputs = run_sidecar(sidecar, fixture_paths, output_dir)

        for spec, fixture_path in zip(FIXTURES, fixture_paths, strict=True):
            source_key = str(fixture_path.resolve())
            parsekit_path = outputs.get(source_key)
            if parsekit_path is None or not parsekit_path.is_file():
                print(f"Missing ParseKit output for {fixture_path}", file=sys.stderr)
                return 1

            baseline_text, baseline_method = BASELINE_EXTRACTORS[spec.file_type](
                fixture_path
            )
            parsekit_text = parsekit_path.read_text(encoding="utf-8", errors="replace")

            baseline_cl100k = count_tokens(baseline_text, "cl100k_base")
            parsekit_cl100k = count_tokens(parsekit_text, "cl100k_base")
            baseline_o200k = count_tokens(baseline_text, "o200k_base")
            parsekit_o200k = count_tokens(parsekit_text, "o200k_base")

            ratio_cl100k = reduction_ratio(baseline_cl100k, parsekit_cl100k)
            ratio_o200k = reduction_ratio(baseline_o200k, parsekit_o200k)

            rows.append(
                {
                    "fixture": spec.filename,
                    "file_type": spec.file_type,
                    "baseline_method": baseline_method,
                    "baseline_cl100k": baseline_cl100k,
                    "parsekit_cl100k": parsekit_cl100k,
                    "delta_cl100k": pct(ratio_cl100k),
                    "baseline_o200k": baseline_o200k,
                    "parsekit_o200k": parsekit_o200k,
                    "delta_o200k": pct(ratio_o200k),
                    "reduction_ratio_cl100k": ratio_cl100k,
                    "why": spec.why,
                }
            )

    write_results(rows, sidecar, generated_at)
    write_ratios(rows, generated_at)

    print(f"Wrote {RESULTS_MD.relative_to(ROOT)}")
    print(f"Wrote {RATIOS_JSON.relative_to(ROOT)}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())